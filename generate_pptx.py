import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen setup
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    COLOR_NAVY = RGBColor(10, 25, 47)      # #0A192F
    COLOR_BG_LIGHT = RGBColor(248, 250, 252) # #F8FAFC
    COLOR_CYAN = RGBColor(6, 182, 212)    # #06B6D4
    COLOR_TEAL = RGBColor(20, 184, 166)   # #14B8A6
    COLOR_DARK = RGBColor(15, 23, 42)     # #0F172A
    COLOR_GRAY = RGBColor(100, 116, 139)  # #64748B
    COLOR_WHITE = RGBColor(255, 255, 255) # #FFFFFF
    COLOR_CARD_BG = RGBColor(255, 255, 255)
    COLOR_GREEN = RGBColor(74, 222, 128)  # #4ADE80

    FONT_HEADING = "Montserrat"
    FONT_BODY = "Arial"
    FONT_CODE = "Consolas"

    def add_slide_transition(slide, transition_type="fade"):
        """Embeds native OpenXML animated slide transitions into PowerPoint."""
        if transition_type == "push":
            xml = r'<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:push dir="l"/></p:transition>'
        elif transition_type == "wipe":
            xml = r'<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:wipe dir="l"/></p:transition>'
        else:
            xml = r'<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>'
        
        transition_elem = parse_xml(xml)
        slide.element.append(transition_elem)

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, subtitle_text):
        # Header title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_HEADING
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK

        # Subtitle
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.name = FONT_BODY
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_TEAL
        p2.font.bold = True

    def set_speaker_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text

    # --- SLIDE 1: COVER ---
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1, COLOR_NAVY)
    add_slide_transition(slide1, "fade")

    # Cover Tag
    shape = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(3.8), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_NAVY
    shape.line.color.rgb = COLOR_CYAN
    shape.line.width = Pt(1)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = "PTI TA 2025/2026 — Executive Summary"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.alignment = PP_ALIGN.CENTER

    # Cover Title
    txBox = slide1.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.5), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PENGANTAR TEKNOLOGI INFORMASI"
    p.font.name = FONT_HEADING
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p2 = tf.add_paragraph()
    p2.text = "Pengenalan Hardware, Sistem Operasi, dan Command Line"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLOR_CYAN

    p3 = tf.add_paragraph()
    p3.text = "Tinjauan Integratif Modul 1, 2, 4, & 8 — Laboratorium Komputer PTI"
    p3.font.name = FONT_BODY
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_GRAY

    # Module Pills
    pills = [
        "Modul 1: PC & Perakitan",
        "Modul 2: Laptop & Perakitan",
        "Modul 4: OS Windows & Linux",
        "Modul 8: Perintah CMD"
    ]
    for idx, pill_text in enumerate(pills):
        x = Inches(0.8 + idx * 2.9)
        shape = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.2), Inches(2.7), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_DARK
        shape.line.color.rgb = COLOR_CYAN
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = pill_text
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER

    set_speaker_notes(slide1, "Selamat pagi/siang rekan-rekan dan bapak/ibu sekalian. Selamat datang dalam presentasi Praktikum Pengantar Teknologi Informasi. Hari ini kita akan membahas rangkuman eksekutif yang mengintegrasikan 4 modul utama: perakitan PC, arsitektur laptop, manajemen OS Windows & Linux, serta perintah Command Line (CMD).")

    # --- SLIDE 2: MODUL 1 (PC HARDWARE) ---
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2, COLOR_BG_LIGHT)
    add_slide_transition(slide2, "push")
    add_header(slide2, "Modul 1: Komponen PC & Perakitan", "Arsitektur Perangkat Keras Komputer Desktop & Prosedur K3 Perakitan")

    # Left Column: 6 Component Cards
    components = [
        ("Motherboard", "Papan induk penghubung & jalur komunikasi utama."),
        ("Processor (CPU)", "Otak pemrosesan instruksi & perhitungan logis."),
        ("RAM Memory", "Penyimpanan sementara berkecepatan tinggi."),
        ("Storage (SSD/HDD)", "Penyimpanan persisten OS, aplikasi, & data."),
        ("VGA Card (GPU)", "Pemrosesan sinyal grafis & akselerasi visual."),
        ("Power Supply (PSU)", "Penyuplai daya listrik DC stabil ke komponen.")
    ]
    for i, (title, desc) in enumerate(components):
        row = i // 2
        col = i % 2
        x = Inches(0.8 + col * 3.3)
        y = Inches(1.5 + row * 1.7)
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.1), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_GRAY
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_DARK
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_GRAY

    # Right Column: Timeline Box
    box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.6), Inches(1.5), Inches(4.9), Inches(5.2))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_WHITE
    box.line.color.rgb = COLOR_CYAN
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Tahapan Perakitan & POST"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_DARK

    steps = [
        "1. Pasang CPU, Thermal Paste, & Cooler",
        "2. Pasang RAM Dual-Channel pada slot utama",
        "3. Instalasi Motherboard ke Casing/Chassis",
        "4. Pasang PSU, Storage, & Cable Management",
        "5. Pengujian POST & Verifikasi Tampilan BIOS"
    ]
    for step in steps:
        p_step = tf.add_paragraph()
        p_step.text = f"\n{step}"
        p_step.font.size = Pt(11)
        p_step.font.color.rgb = COLOR_DARK

    set_speaker_notes(slide2, "Pada Modul 1, kita mempelajari fondasi fisik komputer desktop. Penting bagi praktisi IT untuk paham fungsi spesifik tiap komponen—termasuk PSU yang menyuplai daya DC. Aspek K3 dan ESD harus diperhatikan saat perakitan sebelum pengujian POST.")

    # --- SLIDE 3: MODUL 2 (LAPTOP HARDWARE) ---
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3, COLOR_BG_LIGHT)
    add_slide_transition(slide3, "push")
    add_header(slide3, "Modul 2: Komponen Laptop & Perakitan", "Arsitektur Ringkas, Efisiensi Daya, & Prosedur Disassembly")

    laptop_cards = [
        ("RAM SODIMM", "Form factor memori ekstra ringkas khusus perangkat portabel. Hemat daya dan ruang."),
        ("M.2 NVMe SSD", "Storage bus PCIe kecepatan tinggi berukuran stik tipis tanpa kabel tambahan."),
        ("SoC & Mobile CPU", "Prosesor TDP rendah dengan sistem pendingin Heatpipe dan Fan terintegrasi.")
    ]
    for idx, (title, desc) in enumerate(laptop_cards):
        x = Inches(0.8 + idx * 3.9)
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(3.6), Inches(2.0))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = COLOR_CYAN if idx < 2 else COLOR_GRAY
        card.line.width = Pt(2) if idx < 2 else Pt(1)
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_DARK
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_GRAY

    # Disassembly Banner
    banner = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.9), Inches(11.7), Inches(2.8))
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLOR_NAVY
    banner.line.color.rgb = COLOR_CYAN
    tf = banner.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PROSEDUR DISASSEMBLY & REASSEMBLY LAPTOP\n"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = COLOR_CYAN

    dis_steps = [
        ("Langkah 01", "Buka Bottom Cover & Putuskan arus Baterai utama (Safety First)."),
        ("Langkah 02", "Pelepasan/Upgrading RAM SODIMM & Storage M.2 NVMe."),
        ("Langkah 03", "Pembersihan kipas & pembaruan pasta termal (Cegah Thermal Throttling)."),
        ("Langkah 04", "Perakitan ulang chassis & verifikasi status komponen di BIOS.")
    ]
    for step_num, step_desc in dis_steps:
        p_step = tf.add_paragraph()
        p_step.text = f"• {step_num}: {step_desc}"
        p_step.font.size = Pt(12)
        p_step.font.color.rgb = COLOR_WHITE

    set_speaker_notes(slide3, "Masuk ke Modul 2, kita beralih ke laptop. Perbedaan utama ada pada efisiensi ruang dan daya (RAM SODIMM & SSD M.2 NVMe). Catatan kritis: pemutusan arus baterai utama wajib dilakukan pertama kali sebelum pembongkaran.")

    # --- SLIDE 4: MODUL 4 (OPERATING SYSTEMS) ---
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4, COLOR_BG_LIGHT)
    add_slide_transition(slide4, "push")
    add_header(slide4, "Modul 4: Operating System (Windows & Linux)", "Perbandingan Platform, Manajemen Sumber Daya, & Partisi Boot")

    # Windows Box
    win_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(5.6), Inches(3.2))
    win_box.fill.solid()
    win_box.fill.fore_color.rgb = COLOR_WHITE
    win_box.line.color.rgb = COLOR_CYAN
    win_box.line.width = Pt(2)
    tf = win_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Windows OS"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_DARK
    win_bullets = [
        "Proprietary & Komersial",
        "Antarmuka GUI sangat User-Friendly",
        "Dukungan Ekosistem Software & Game Luas",
        "Standar Lingkungan Kerja & Perkantoran"
    ]
    for b in win_bullets:
        pb = tf.add_paragraph()
        pb.text = f"• {b}"
        pb.font.size = Pt(11)
        pb.font.color.rgb = COLOR_GRAY

    # Linux Box
    lin_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(5.6), Inches(3.2))
    lin_box.fill.solid()
    lin_box.fill.fore_color.rgb = COLOR_WHITE
    lin_box.line.color.rgb = RGBColor(249, 115, 22) # Orange
    lin_box.line.width = Pt(2)
    tf = lin_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Linux (Ubuntu)"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_DARK
    lin_bullets = [
        "Open-Source & Gratis",
        "Sangat Aman & Relatif Bebas Virus",
        "Stabilitas Tinggi untuk Server & Developer",
        "Manajemen Partisi Fleksibel (ext4 & Swap)"
    ]
    for b in lin_bullets:
        pb = tf.add_paragraph()
        pb.text = f"• {b}"
        pb.font.size = Pt(11)
        pb.font.color.rgb = COLOR_GRAY

    # 4 Functions Grid
    funcs = [
        ("Manajemen Proses", "CPU Scheduling & RAM"),
        ("File System", "Struktur NTFS / ext4"),
        ("Abstraksi Hardware", "Driver & I/O Management"),
        ("User Interface", "Antarmuka GUI & CLI")
    ]
    for idx, (title, desc) in enumerate(funcs):
        x = Inches(0.8 + idx * 2.95)
        fcard = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.9), Inches(2.8), Inches(1.8))
        fcard.fill.solid()
        fcard.fill.fore_color.rgb = COLOR_WHITE
        fcard.line.color.rgb = COLOR_TEAL
        tf = fcard.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_DARK
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_GRAY

    set_speaker_notes(slide4, "Pada Modul 4, kita membahas OS sebagai pengelola sumber daya. Kita membandingkan Windows yang unggul di antarmuka GUI, dengan Linux Ubuntu yang menjadi standar infrastruktur server karena ringan, open-source, dan aman.")

    # --- SLIDE 5: MODUL 8 (COMMAND PROMPT) ---
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5, COLOR_BG_LIGHT)
    add_slide_transition(slide5, "push")
    add_header(slide5, "Modul 8: Command Line Interface (CMD)", "Eksekusi Perintah Teks, Automasi, & Diagnostics")

    # Terminal Window Box
    term = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2))
    term.fill.solid()
    term.fill.fore_color.rgb = RGBColor(5, 14, 26)
    term.line.color.rgb = COLOR_CYAN
    tf = term.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Windows Command Prompt [Administrator]\n"
    p.font.name = FONT_CODE
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_CYAN

    cmd_data = [
        ("DIR", "Melihat daftar file & direktori", "dir /w"),
        ("CD", "Berpindah direktori aktif", "cd Documents  /  cd .."),
        ("MD / MKDIR", "Membuat folder baru", "mkdir Project_PTI"),
        ("RD / RMDIR", "Menghapus folder", "rd /s Old_Folder"),
        ("CLS", "Membersihkan layar terminal", "cls"),
        ("IPCONFIG", "Cek IP & konfigurasi jaringan", "ipconfig /all")
    ]
    for cmd, desc, ex in cmd_data:
        p_cmd = tf.add_paragraph()
        p_cmd.text = f"  {cmd:<12} |  {desc:<35} |  Ex: {ex}"
        p_cmd.font.name = FONT_CODE
        p_cmd.font.size = Pt(11)
        p_cmd.font.color.rgb = COLOR_GREEN

    set_speaker_notes(slide5, "Di Modul 8, kita memanfaatkan antarmuka CLI berbasis teks. Menguasai CMD seperti ipconfig, dir, dan mkdir merupakan batu loncatan penting sebelum mempelajari PowerShell, Linux Shell, dan automation scripting.")

    # --- SLIDE 6: KESIMPULAN ---
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide6, COLOR_BG_LIGHT)
    add_slide_transition(slide6, "fade")
    add_header(slide6, "Integrasi Modul Praktikum PTI", "Sinergi Layer Komputasi: Hardware, OS, & Command Line")

    layers = [
        ("1. HARDWARE", "Modul 1 & 2: Fondasi fisik komputasi (PC & Laptop)"),
        ("2. OPERATING SYSTEM", "Modul 4: Pengelola daya komputasi & memori"),
        ("3. COMMAND LINE", "Modul 8: Kendali presisi & automasi sistem")
    ]
    for idx, (title, desc) in enumerate(layers):
        x = Inches(0.8 + idx * 3.9)
        lbox = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(3.6), Inches(2.0))
        lbox.fill.solid()
        lbox.fill.fore_color.rgb = COLOR_WHITE
        lbox.line.color.rgb = COLOR_TEAL
        lbox.line.width = Pt(2)
        tf = lbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_DARK
        p2 = tf.add_paragraph()
        p2.text = f"\n{desc}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_GRAY

    # Executive Summary Box
    summary = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.0), Inches(11.7), Inches(2.6))
    summary.fill.solid()
    summary.fill.fore_color.rgb = COLOR_NAVY
    summary.line.color.rgb = COLOR_CYAN
    summary.line.width = Pt(2)
    tf = summary.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "KESIMPULAN EKSEKUTIF"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_CYAN
    p.alignment = PP_ALIGN.CENTER

    p_quote = tf.add_paragraph()
    p_quote.text = "\n\"Pemahaman IT yang utuh tidak sekadar mampu menggunakan komputer, melainkan menguasai bagaimana komputer dirakit, dijalankan sistemnya, dan dikendalikan secara efisien.\""
    p_quote.font.size = Pt(15)
    p_quote.font.color.rgb = COLOR_WHITE
    p_quote.font.italic = True
    p_quote.alignment = PP_ALIGN.CENTER

    set_speaker_notes(slide6, "Sebagai penutup, keempat modul ini membentuk satu kesatuan siklus komputasi yang utuh: Hardware adalah fondasinya, OS adalah pengelolanya, dan CLI adalah alat kendali presisinya. Terima kasih!")

    output_filename = "presentasi_pti.pptx"
    prs.save(output_filename)
    print(f"File PowerPoint berhasil dibuat: {os.path.abspath(output_filename)}")

if __name__ == "__main__":
    create_presentation()